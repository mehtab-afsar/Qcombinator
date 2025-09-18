"""
Document Upload & Processing Module (DUP)
Extracts structured data from pitch decks, financials, and business plans
"""

import os
import io
from typing import Dict, Any
from datetime import datetime
import pypdf
from pptx import Presentation
import pandas as pd
from fastapi import UploadFile

class DocumentProcessor:
    """
    Process uploaded documents and extract structured data
    Supports: PDF, PPT/PPTX, XLS/XLSX
    """
    
    async def process(self, file: UploadFile) -> Dict[str, Any]:
        """
        Main processing function
        Returns extracted data with page references
        """
        
        # Read file content
        content = await file.read()
        file_ext = os.path.splitext(file.filename)[1].lower()
        
        extracted_data = {
            'filename': file.filename,
            'file_type': file_ext,
            'processed_at': datetime.utcnow().isoformat(),
            'sections': {},
            'metadata': {}
        }
        
        # Process based on file type
        if file_ext == '.pdf':
            extracted_data.update(self._process_pdf(content))
        elif file_ext in ['.ppt', '.pptx']:
            extracted_data.update(self._process_pptx(content))
        elif file_ext in ['.xls', '.xlsx']:
            extracted_data.update(self._process_excel(content))
        
        # Extract key information
        extracted_data['company_info'] = self._extract_company_info(extracted_data)
        extracted_data['financial_data'] = self._extract_financial_data(extracted_data)
        extracted_data['team_data'] = self._extract_team_data(extracted_data)
        extracted_data['market_data'] = self._extract_market_data(extracted_data)
        
        return extracted_data
    
    def _process_pdf(self, content: bytes) -> Dict:
        """Extract text and structure from PDF"""
        pdf_file = io.BytesIO(content)
        pdf_reader = pypdf.PdfReader(pdf_file)
        
        pages = []
        full_text = ""
        
        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            pages.append({
                'page': page_num + 1,
                'content': text
            })
            full_text += text + "\n"
        
        return {
            'pages': pages,
            'full_text': full_text,
            'page_count': len(pages)
        }
    
    def _process_pptx(self, content: bytes) -> Dict:
        """Extract text and structure from PowerPoint"""
        pptx_file = io.BytesIO(content)
        presentation = Presentation(pptx_file)
        
        slides = []
        full_text = ""
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = ""
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            
            slides.append({
                'slide': slide_num + 1,
                'title': slide.shapes.title.text if slide.shapes.title else "",
                'content': slide_text
            })
            full_text += slide_text + "\n"
        
        return {
            'slides': slides,
            'full_text': full_text,
            'slide_count': len(slides)
        }
    
    def _process_excel(self, content: bytes) -> Dict:
        """Extract data from Excel files"""
        excel_file = io.BytesIO(content)
        
        sheets = {}
        try:
            # Read all sheets
            excel_data = pd.read_excel(excel_file, sheet_name=None)
            
            for sheet_name, df in excel_data.items():
                sheets[sheet_name] = {
                    'data': df.to_dict('records'),
                    'columns': df.columns.tolist(),
                    'row_count': len(df)
                }
        except Exception as e:
            sheets['error'] = str(e)
        
        return {
            'sheets': sheets,
            'sheet_count': len(sheets)
        }
    
    def _extract_company_info(self, data: Dict) -> Dict:
        """Extract company information from document"""
        info = {
            'name': None,
            'description': None,
            'industry': None,
            'founded': None,
            'website': None,
            'location': None
        }
        
        # Simple extraction logic (would use NLP/LLM in production)
        text = data.get('full_text', '')
        
        # Look for common patterns
        lines = text.split('\n')
        for line in lines[:50]:  # Check first 50 lines
            lower_line = line.lower()
            
            if 'founded' in lower_line or 'established' in lower_line:
                info['founded'] = line
            elif 'industry' in lower_line or 'sector' in lower_line:
                info['industry'] = line
            elif 'www.' in lower_line or 'http' in lower_line:
                info['website'] = line
            elif any(loc in lower_line for loc in ['san francisco', 'new york', 'boston', 'austin']):
                info['location'] = line
        
        return info
    
    def _extract_financial_data(self, data: Dict) -> Dict:
        """Extract financial metrics from document"""
        financials = {
            'revenue': None,
            'mrr': None,
            'arr': None,
            'burn_rate': None,
            'runway': None,
            'gross_margin': None,
            'cac': None,
            'ltv': None,
            'growth_rate': None
        }
        
        # Check if Excel data contains financials
        if 'sheets' in data:
            for sheet_name, sheet_data in data.get('sheets', {}).items():
                if 'financial' in sheet_name.lower() or 'revenue' in sheet_name.lower():
                    # Extract from structured data
                    if 'data' in sheet_data and sheet_data['data']:
                        # Simple extraction (would be more sophisticated in production)
                        for row in sheet_data['data']:
                            for key, value in row.items():
                                if 'revenue' in str(key).lower():
                                    financials['revenue'] = value
                                elif 'mrr' in str(key).lower():
                                    financials['mrr'] = value
                                elif 'burn' in str(key).lower():
                                    financials['burn_rate'] = value
        
        # Also check text for financial mentions
        text = data.get('full_text', '')
        
        # Look for common financial patterns (simplified)
        import re
        
        # Revenue pattern
        revenue_pattern = r'\$?\d+\.?\d*[MmKk]?\s*(?:revenue|ARR|MRR)'
        revenue_matches = re.findall(revenue_pattern, text)
        if revenue_matches and not financials['revenue']:
            financials['revenue'] = revenue_matches[0]
        
        return financials
    
    def _extract_team_data(self, data: Dict) -> Dict:
        """Extract team information from document"""
        team = {
            'founders': [],
            'team_size': None,
            'key_members': []
        }
        
        text = data.get('full_text', '')
        
        # Look for team slide/section
        if 'slides' in data:
            for slide in data['slides']:
                if 'team' in slide.get('title', '').lower():
                    # Extract names (simplified - would use NER in production)
                    content = slide.get('content', '')
                    lines = content.split('\n')
                    for line in lines:
                        if any(title in line.lower() for title in ['ceo', 'cto', 'founder', 'co-founder']):
                            team['key_members'].append(line.strip())
        
        return team
    
    def _extract_market_data(self, data: Dict) -> Dict:
        """Extract market information from document"""
        market = {
            'tam': None,
            'sam': None,
            'som': None,
            'competitors': [],
            'growth_rate': None
        }
        
        text = data.get('full_text', '')
        
        # Look for TAM/SAM/SOM mentions
        import re
        
        tam_pattern = r'TAM.*?\$?\d+\.?\d*[BbMm]'
        tam_matches = re.findall(tam_pattern, text)
        if tam_matches:
            market['tam'] = tam_matches[0]
        
        sam_pattern = r'SAM.*?\$?\d+\.?\d*[BbMm]'
        sam_matches = re.findall(sam_pattern, text)
        if sam_matches:
            market['sam'] = sam_matches[0]
        
        return market